"""Enterprise Agent OS — Document Parser.

Parses multiple document formats (PDF, DOCX, PPTX, XLSX, CSV, TXT)
into structured content with text, tables, and metadata extraction.

Supports:
- PDF: text extraction via pypdf
- DOCX: text + table extraction via python-docx
- PPTX: slide text extraction via python-pptx
- XLSX: sheet data extraction via openpyxl
- CSV: tabular data parsing
- TXT/MD: plain text with encoding detection
"""
from __future__ import annotations
import csv
import io
import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple


@dataclass
class ParsedDocument:
    """Result of parsing a document."""
    filename: str
    content: str
    doc_type: str
    metadata: Dict[str, Any] = field(default_factory=dict)
    tables: List[Dict[str, Any]] = field(default_factory=list)
    images: List[Dict[str, Any]] = field(default_factory=list)
    sections: List[Dict[str, Any]] = field(default_factory=list)

    @property
    def text_length(self) -> int:
        """Get total text length."""
        return len(self.content)

    @property
    def table_count(self) -> int:
        """Get number of tables."""
        return len(self.tables)

    @property
    def summary(self) -> str:
        """Get a brief summary of the parsed document."""
        return (
            f"{self.filename}: {self.doc_type.upper()} | "
            f"{self.text_length} chars | "
            f"{self.table_count} tables | "
            f"{len(self.sections)} sections"
        )


def parse_pdf(file_path: Path) -> ParsedDocument:
    """Parse a PDF file and extract text content.

    Args:
        file_path: Path to the PDF file.

    Returns:
        ParsedDocument with extracted text and metadata.
    """
    text = ""
    metadata: Dict[str, Any] = {"pages": 0}

    try:
        from pypdf import PdfReader
        reader = PdfReader(str(file_path))
        metadata["pages"] = len(reader.pages)

        # Extract metadata if available
        if reader.metadata:
            metadata["title"] = getattr(reader.metadata, "title", None)
            metadata["author"] = getattr(reader.metadata, "author", None)

        pages = []
        for i, page in enumerate(reader.pages):
            page_text = page.extract_text() or ""
            pages.append({"page": i + 1, "text": page_text})

        text = "\n\n".join(p["text"] for p in pages)
        metadata["page_texts"] = pages

    except ImportError:
        text = f"[PDF file: {file_path.name}] — pypdf not installed"
    except Exception as e:
        text = f"[PDF parsing error: {e}]"

    return ParsedDocument(
        filename=file_path.name,
        content=text,
        doc_type="pdf",
        metadata=metadata,
    )


def parse_docx(file_path: Path) -> ParsedDocument:
    """Parse a DOCX file and extract text + tables.

    Args:
        file_path: Path to the DOCX file.

    Returns:
        ParsedDocument with text, tables, and sections.
    """
    text = ""
    tables: List[Dict[str, Any]] = []
    sections: List[Dict[str, Any]] = []
    metadata: Dict[str, Any] = {}

    try:
        from docx import Document
        doc = Document(str(file_path))

        # Extract paragraphs with heading detection
        paragraphs = []
        for para in doc.paragraphs:
            if para.style and para.style.name.startswith("Heading"):
                level = para.style.name.replace("Heading", "").strip() or "1"
                sections.append({
                    "title": para.text,
                    "level": int(level) if level.isdigit() else 1,
                    "position": len(paragraphs),
                })
            paragraphs.append(para.text)

        text = "\n".join(paragraphs)

        # Extract tables
        for i, table in enumerate(doc.tables):
            headers = []
            rows = []
            for j, row in enumerate(table.rows):
                cells = [cell.text.strip() for cell in row.cells]
                if j == 0:
                    headers = cells
                else:
                    rows.append(cells)
            tables.append({
                "id": f"docx_table_{i}",
                "headers": headers,
                "rows": rows,
                "caption": f"Table {i + 1}",
            })

        metadata["paragraph_count"] = len(paragraphs)
        metadata["table_count"] = len(tables)

    except ImportError:
        text = f"[DOCX file: {file_path.name}] — python-docx not installed"
    except Exception as e:
        text = f"[DOCX parsing error: {e}]"

    return ParsedDocument(
        filename=file_path.name,
        content=text,
        doc_type="docx",
        metadata=metadata,
        tables=tables,
        sections=sections,
    )


def parse_pptx(file_path: Path) -> ParsedDocument:
    """Parse a PPTX file and extract slide text.

    Args:
        file_path: Path to the PPTX file.

    Returns:
        ParsedDocument with slide text and metadata.
    """
    text = ""
    sections: List[Dict[str, Any]] = []
    metadata: Dict[str, Any] = {}

    try:
        from pptx import Presentation
        prs = Presentation(str(file_path))

        slides_text = []
        for i, slide in enumerate(prs.slides):
            slide_content = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content.append(shape.text.strip())

            slide_text = "\n".join(slide_content)
            slides_text.append(f"--- Slide {i + 1} ---\n{slide_text}")

            sections.append({
                "title": f"Slide {i + 1}",
                "text": slide_text,
                "position": i,
            })

        text = "\n\n".join(slides_text)
        metadata["slide_count"] = len(prs.slides)

    except ImportError:
        text = f"[PPTX file: {file_path.name}] — python-pptx not installed"
    except Exception as e:
        text = f"[PPTX parsing error: {e}]"

    return ParsedDocument(
        filename=file_path.name,
        content=text,
        doc_type="pptx",
        metadata=metadata,
        sections=sections,
    )


def parse_xlsx(file_path: Path) -> ParsedDocument:
    """Parse an XLSX file and extract sheet data as tables.

    Args:
        file_path: Path to the XLSX file.

    Returns:
        ParsedDocument with table data from each sheet.
    """
    text = ""
    tables: List[Dict[str, Any]] = []
    metadata: Dict[str, Any] = {}

    try:
        from openpyxl import load_workbook
        wb = load_workbook(str(file_path), read_only=True, data_only=True)

        metadata["sheet_names"] = wb.sheetnames
        sheets_text = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            sheet_data = []
            headers = []
            rows = []

            for i, row in enumerate(ws.iter_rows(values_only=True)):
                cells = [str(cell) if cell is not None else "" for cell in row]
                if i == 0:
                    headers = cells
                else:
                    rows.append(cells)
                sheet_data.append(cells)

            # Convert to text
            sheet_text_lines = [f"Sheet: {sheet_name}"]
            if headers:
                sheet_text_lines.append(" | ".join(headers))
            for row in rows:
                sheet_text_lines.append(" | ".join(row))

            sheets_text.append("\n".join(sheet_text_lines))

            tables.append({
                "id": f"xlsx_{sheet_name}",
                "headers": headers,
                "rows": rows[:100],  # Limit rows for memory
                "caption": f"Sheet: {sheet_name}",
                "total_rows": len(rows),
            })

        text = "\n\n".join(sheets_text)
        metadata["table_count"] = len(tables)
        wb.close()

    except ImportError:
        text = f"[XLSX file: {file_path.name}] — openpyxl not installed"
    except Exception as e:
        text = f"[XLSX parsing error: {e}]"

    return ParsedDocument(
        filename=file_path.name,
        content=text,
        doc_type="xlsx",
        metadata=metadata,
        tables=tables,
    )


def parse_csv(file_path: Path) -> ParsedDocument:
    """Parse a CSV file into structured table data.

    Args:
        file_path: Path to the CSV file.

    Returns:
        ParsedDocument with CSV data as table.
    """
    text = ""
    tables: List[Dict[str, Any]] = []
    metadata: Dict[str, Any] = {}

    try:
        raw = file_path.read_text(encoding="utf-8", errors="replace")
        reader = csv.reader(io.StringIO(raw))
        all_rows = list(reader)

        if not all_rows:
            return ParsedDocument(
                filename=file_path.name,
                content="",
                doc_type="csv",
                metadata={"rows": 0},
            )

        headers = all_rows[0]
        rows = all_rows[1:]

        # Build text representation
        text_lines = [" | ".join(headers)]
        for row in rows:
            text_lines.append(" | ".join(row))
        text = "\n".join(text_lines)

        tables.append({
            "id": f"csv_{file_path.stem}",
            "headers": headers,
            "rows": rows[:100],  # Limit for memory
            "caption": file_path.stem,
            "total_rows": len(rows),
        })

        metadata["rows"] = len(rows)
        metadata["columns"] = len(headers)

    except Exception as e:
        text = f"[CSV parsing error: {e}]"

    return ParsedDocument(
        filename=file_path.name,
        content=text,
        doc_type="csv",
        metadata=metadata,
        tables=tables,
    )


def parse_text(file_path: Path) -> ParsedDocument:
    """Parse a plain text or markdown file.

    Args:
        file_path: Path to the text file.

    Returns:
        ParsedDocument with text content.
    """
    try:
        text = file_path.read_text(encoding="utf-8", errors="replace")
    except Exception as e:
        text = f"[Text parsing error: {e}]"

    # Detect if markdown
    is_markdown = file_path.suffix.lower() in (".md", ".markdown")
    doc_type = "markdown" if is_markdown else "text"

    # Extract sections from markdown headings
    sections = []
    if is_markdown:
        for match in re.finditer(r"^(#{1,6})\s+(.+)$", text, re.MULTILINE):
            level = len(match.group(1))
            sections.append({
                "title": match.group(2),
                "level": level,
                "position": match.start(),
            })

    return ParsedDocument(
        filename=file_path.name,
        content=text,
        doc_type=doc_type,
        metadata={"encoding": "utf-8"},
        sections=sections,
    )


# Extension to parser mapping
_PARSER_MAP = {
    ".pdf": parse_pdf,
    ".docx": parse_docx,
    ".pptx": parse_pptx,
    ".xlsx": parse_xlsx,
    ".xls": parse_xlsx,
    ".csv": parse_csv,
    ".tsv": parse_csv,
    ".txt": parse_text,
    ".md": parse_text,
    ".markdown": parse_text,
    ".log": parse_text,
    ".json": parse_text,
    ".xml": parse_text,
    ".html": parse_text,
    ".htm": parse_text,
}


def parse_document(file_path: str | Path) -> ParsedDocument:
    """Auto-detect format and parse a document.

    Args:
        file_path: Path to the document file.

    Returns:
        ParsedDocument with extracted content.
    """
    path = Path(file_path)
    if not path.exists():
        return ParsedDocument(
            filename=path.name,
            content=f"[File not found: {path}]",
            doc_type="error",
            metadata={"error": "file_not_found"},
        )

    suffix = path.suffix.lower()
    parser_func = _PARSER_MAP.get(suffix)

    if parser_func:
        return parser_func(path)

    # Fallback: try as text
    return parse_text(path)
